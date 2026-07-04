"""
Extracts raw text (and simple metadata) from uploaded PDF, DOCX, PPTX, and
XLSX files, then splits it into overlapping chunks for embedding.
"""
from dataclasses import dataclass
from typing import List

from app.utils.logger import get_logger

logger = get_logger(__name__)


@dataclass
class ExtractedPage:
    page_number: int
    text: str


@dataclass
class ExtractionResult:
    pages: List[ExtractedPage]
    page_count: int


def extract_text(file_path: str, file_type: str) -> ExtractionResult:
    file_type = file_type.lower()
    if file_type == "pdf":
        return _extract_pdf(file_path)
    if file_type in ("docx", "doc"):
        return _extract_docx(file_path)
    if file_type in ("pptx", "ppt"):
        return _extract_pptx(file_path)
    if file_type in ("xlsx", "xls", "csv"):
        return _extract_xlsx(file_path)
    raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf(file_path: str) -> ExtractionResult:
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    pages = [
        ExtractedPage(page_number=i + 1, text=page.extract_text() or "")
        for i, page in enumerate(reader.pages)
    ]
    return ExtractionResult(pages=pages, page_count=len(pages))


def _extract_docx(file_path: str) -> ExtractionResult:
    import docx

    document = docx.Document(file_path)
    text = "\n".join(p.text for p in document.paragraphs)
    return ExtractionResult(pages=[ExtractedPage(1, text)], page_count=1)


def _extract_pptx(file_path: str) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text for shape in slide.shapes if shape.has_text_frame
        ]
        pages.append(ExtractedPage(page_number=i + 1, text="\n".join(texts)))
    return ExtractionResult(pages=pages, page_count=len(pages))


def _extract_xlsx(file_path: str) -> ExtractionResult:
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    pages = []
    for i, sheet in enumerate(wb.worksheets):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(" | ".join(str(c) for c in row if c is not None))
        pages.append(ExtractedPage(page_number=i + 1, text="\n".join(rows)))
    return ExtractionResult(pages=pages, page_count=len(pages))


def chunk_text(pages: List[ExtractedPage], chunk_size: int = 1000, overlap: int = 150) -> List[dict]:
    """Splits page text into overlapping chunks, tracking source page for citations."""
    chunks = []
    idx = 0
    for page in pages:
        text = page.text.strip()
        if not text:
            continue
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunk = text[start:end]
            chunks.append({
                "chunk_index": idx,
                "page_number": page.page_number,
                "text": chunk,
            })
            idx += 1
            if end == len(text):
                break
            start = end - overlap
    return chunks
